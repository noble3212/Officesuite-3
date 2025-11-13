# === Imports ===
import os
from tkinter import *
from tkinter import ttk, filedialog, messagebox, font  # Added font
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader, PdfWriter

# === Office Suite App ===
class OfficeSuiteApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Local Office Suite")
        self.root.geometry("400x400")

        Label(root, text="Local Office Suite", font=("Arial", 18)).pack(pady=20)

        Button(root, text="Word", width=20, command=self.open_word_window).pack(pady=5)
        Button(root, text="Excel", width=20, command=self.open_excel_window).pack(pady=5)
        Button(root, text="PowerPoint", width=20, command=self.open_powerpoint_window).pack(pady=5)
        Button(root, text="PDF Merge", width=20, command=self.open_pdf_window).pack(pady=5)

    # === Word Editor ===
    def open_word_window(self):
        win = Toplevel(self.root)
        win.title("Word Editor")
        win.geometry("600x500")

        Label(win, text="Word Editor", font=("Arial",16)).pack(pady=5)

        # --- Font Controls Frame ---
        font_frame = Frame(win)
        font_frame.pack(pady=5)

        # Get available fonts
        font_families = font.families()
        self.word_font_family = StringVar(win)
        self.word_font_family.set(font_families[0]) # Default
        
        font_family_menu = ttk.Combobox(font_frame, textvariable=self.word_font_family, values=font_families, state='readonly', width=25)
        font_family_menu.pack(side=LEFT, padx=5)

        font_sizes = [8, 9, 10, 11, 12, 14, 16, 18, 24, 32]
        self.word_font_size = IntVar(win)
        self.word_font_size.set(12) # Default

        font_size_menu = ttk.Combobox(font_frame, textvariable=self.word_font_size, values=font_sizes, width=5)
        font_size_menu.pack(side=LEFT, padx=5)
        
        # --- Style Buttons ---
        bold_btn = Button(font_frame, text="B", width=3, command=self.toggle_bold)
        bold_btn.pack(side=LEFT, padx=1)
        
        italic_btn = Button(font_frame, text="I", width=3, command=self.toggle_italic)
        italic_btn.pack(side=LEFT, padx=1)
        
        underline_btn = Button(font_frame, text="U", width=3, command=self.toggle_underline)
        underline_btn.pack(side=LEFT, padx=1)


        self.text_area = Text(win, wrap=WORD)
        self.text_area.pack(fill=BOTH, expand=True, padx=10, pady=10)

        def apply_font(event=None):
            try:
                family = self.word_font_family.get()
                size = self.word_font_size.get()
                
                # Base font
                base_font = font.Font(family=family, size=size)
                self.text_area.configure(font=base_font)
                
                # Configure style fonts and tags
                self.bold_font = font.Font(family=family, size=size, weight="bold")
                self.text_area.tag_configure("bold", font=self.bold_font)
                
                self.italic_font = font.Font(family=family, size=size, slant="italic")
                self.text_area.tag_configure("italic", font=self.italic_font)
                
                self.underline_font = font.Font(family=family, size=size, underline=True)
                self.text_area.tag_configure("underline", font=self.underline_font)
                
            except Exception as e:
                messagebox.showerror("Font Error", f"Could not apply font: {e}")

        font_family_menu.bind("<<ComboboxSelected>>", apply_font)
        font_size_menu.bind("<<ComboboxSelected>>", apply_font)
        
        # Initialize fonts and tags
        apply_font()

        def save_word():
            # Note: This simple save does not preserve font formatting.
            filename = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word Doc","*.docx")])
            if filename:
                doc = Document()
                doc.add_paragraph(self.text_area.get("1.0", END))
                doc.save(filename)
                messagebox.showinfo("Saved", f"Word document saved: {filename}")

        def upload_word():
            filename = filedialog.askopenfilename(filetypes=[("Word Doc","*.docx")])
            if filename:
                doc = Document(filename)
                self.text_area.delete("1.0", END)
                for para in doc.paragraphs:
                    self.text_area.insert(END, para.text + "\n")
        
        # --- Button Frame ---
        btn_frame_word = Frame(win)
        btn_frame_word.pack(pady=5)
        Button(btn_frame_word, text="Upload Word", command=upload_word).pack(side=LEFT, padx=5, pady=5)
        Button(btn_frame_word, text="Save Word", command=save_word).pack(side=LEFT, padx=5, pady=5)

    # --- Text Style Toggle Functions ---
    def toggle_style(self, tag_name):
        try:
            # Check if tag is applied
            current_tags = self.text_area.tag_names("sel.first")
            if tag_name in current_tags:
                self.text_area.tag_remove(tag_name, "sel.first", "sel.last")
            else:
                self.text_area.tag_add(tag_name, "sel.first", "sel.last")
        except TclError:
            # No text selected
            pass
            
    def toggle_bold(self):
        self.toggle_style("bold")

    def toggle_italic(self):
        self.toggle_style("italic")

    def toggle_underline(self):
        self.toggle_style("underline")

    # === Excel Editor ===
    def open_excel_window(self):
        win = Toplevel(self.root)
        win.title("Excel Editor")
        win.geometry("900x500")

        Label(win, text="Excel Editor", font=("Arial",16)).pack(pady=5)
        
        # --- Expanded columns ---
        columns = tuple(chr(ord('A') + i) for i in range(10)) # A-J
        num_cols = len(columns)

        tree = ttk.Treeview(win, columns=columns, show="headings", selectmode="extended")
        for col in columns:
            tree.heading(col, text=col)
            tree.column(col, width=100, anchor=CENTER)
        tree.pack(fill=BOTH, expand=True, padx=10, pady=10)

        # --- Cell editing ---
        def on_double_click(event):
            item = tree.identify('item', event.x, event.y)
            column = tree.identify_column(event.x)
            if not item or not column:
                return
            x, y, width, height = tree.bbox(item, column)
            value = tree.set(item, column)
            entry = Entry(tree)
            entry.place(x=x, y=y, width=width, height=height)
            entry.insert(0, value)
            entry.focus()

            def save_edit(e):
                tree.set(item, column, entry.get())
                entry.destroy()

            entry.bind("<Return>", save_edit)
            entry.bind("<FocusOut>", save_edit)

        tree.bind("<Double-1>", on_double_click)

        # --- Add initial empty rows (10 rows) ---
        for _ in range(10):
            tree.insert("", END, values=[""]*num_cols)

        # --- Excel Functions ---
        def add_row():
            tree.insert("", END, values=[""]*num_cols)

        def get_selected_values():
            selected = tree.selection()
            values = []
            for item in selected:
                row = tree.item(item)['values']
                for val in row:
                    try:
                        values.append(float(val))
                    except:
                        continue
            return values

        def sum_selected(event=None):
            vals = get_selected_values()
            messagebox.showinfo("Sum", f"Sum: {sum(vals)}")

        def avg_selected(event=None):
            vals = get_selected_values()
            if vals:
                messagebox.showinfo("Average", f"Average: {sum(vals)/len(vals)}")
            else:
                messagebox.showinfo("Average", "No numeric values selected")

        def min_selected(event=None):
            vals = get_selected_values()
            if vals:
                messagebox.showinfo("Min", f"Min: {min(vals)}")
            else:
                messagebox.showinfo("Min", "No numeric values selected")

        def max_selected(event=None):
            vals = get_selected_values()
            if vals:
                messagebox.showinfo("Max", f"Max: {max(vals)}")
            else:
                messagebox.showinfo("Max", "No numeric values selected")

        def save_excel():
            filename = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel File","*.xlsx")])
            if filename:
                wb = Workbook()
                ws = wb.active
                # Add headers
                ws.append(list(columns))
                # Add data
                for i in tree.get_children():
                    ws.append(tree.item(i)['values'])
                wb.save(filename)
                messagebox.showinfo("Saved", f"Excel file saved: {filename}")

        def upload_excel():
            filename = filedialog.askopenfilename(filetypes=[("Excel File","*.xlsx")])
            if filename:
                wb = load_workbook(filename)
                ws = wb.active
                tree.delete(*tree.get_children())
                
                # Check for headers
                header_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), None)
                
                # Update columns in treeview if necessary (simple check)
                if header_row and len(header_row) > num_cols:
                    new_cols = tuple(str(h) for h in header_row)
                    tree["columns"] = new_cols
                    for col in new_cols:
                         tree.heading(col, text=col)
                         tree.column(col, width=100, anchor=CENTER)
                
                # Load data
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i == 0 and list(row) == list(header_row): # Skip header if it matches
                        continue
                    tree.insert("", END, values=[str(c) if c is not None else "" for c in row])


        # --- Buttons ---
        btn_frame = Frame(win)
        btn_frame.pack()
        Button(btn_frame, text="Add Row", command=add_row).pack(side=LEFT, pady=2, padx=5)
        Button(btn_frame, text="Upload Excel", command=upload_excel).pack(side=LEFT, pady=2, padx=5)
        Button(btn_frame, text="Save Excel", command=save_excel).pack(side=LEFT, pady=2, padx=5)

        # --- Keybinds ---
        win.bind("<Control-m>", sum_selected)
        win.bind("<Control-a>", avg_selected)
        win.bind("<Control-x>", min_selected)
        win.bind("<Control-z>", max_selected)

    # === PowerPoint Editor ===
    def open_powerpoint_window(self):
        win = Toplevel(self.root)
        win.title("PowerPoint Editor")
        win.geometry("800x500")
        
        Label(win, text="PowerPoint Editor", font=("Arial",16)).pack(pady=5)
        
        # --- Main Layout ---
        # *** FIX: Changed PanedWindow to ttk.PanedWindow to fix crash ***
        main_pane = ttk.PanedWindow(win, orient=HORIZONTAL)
        main_pane.pack(fill=BOTH, expand=True, padx=10, pady=10)

        # --- Left Frame (Slide List) ---
        left_frame = ttk.Frame(main_pane, padding=5)
        Label(left_frame, text="Slides").pack()
        self.ppt_slide_listbox = Listbox(left_frame, exportselection=False)
        self.ppt_slide_listbox.pack(fill=BOTH, expand=True)
        main_pane.add(left_frame, weight=1) # This now works

        # --- Right Frame (Editor) ---
        right_frame = ttk.Frame(main_pane, padding=5)
        Label(right_frame, text="Slide Title:").pack(pady=5)
        self.ppt_slide_title = Entry(right_frame, width=50)
        self.ppt_slide_title.pack(pady=5, fill=X, padx=5)
        
        Label(right_frame, text="Slide Content:").pack(pady=5)
        self.ppt_slide_content = Text(right_frame, height=10, wrap=WORD)
        self.ppt_slide_content.pack(padx=5, pady=5, fill=BOTH, expand=True)
        main_pane.add(right_frame, weight=3) # This now works
        
        self.ppt_slides = [] # Holds the slide data

        def clear_ppt_fields():
            self.ppt_slide_title.delete(0, END)
            self.ppt_slide_content.delete("1.0", END)
            self.ppt_slide_listbox.selection_clear(0, END)

        def add_slide():
            title = self.ppt_slide_title.get() or "Untitled"
            content = self.ppt_slide_content.get("1.lock", END)
            
            self.ppt_slides.append({"title": title, "content": content})
            self.ppt_slide_listbox.insert(END, f"{len(self.ppt_slides)}: {title}")
            
            clear_ppt_fields()
            messagebox.showinfo("Slide Added", "Slide added to presentation")
        
        def save_slide_edit():
            selected_indices = self.ppt_slide_listbox.curselection()
            if not selected_indices:
                messagebox.showwarning("No Selection", "Please select a slide from the list to save edits.")
                return
            
            index = selected_indices[0]
            new_title = self.ppt_slide_title.get() or "Untitled"
            new_content = self.ppt_slide_content.get("1.0", END)
            
            # Update data
            self.ppt_slides[index] = {"title": new_title, "content": new_content}
            
            # Update listbox
            self.ppt_slide_listbox.delete(index)
            self.ppt_slide_listbox.insert(index, f"{index + 1}: {new_title}")
            self.ppt_slide_listbox.selection_set(index)
            messagebox.showinfo("Slide Saved", "Your edits to the selected slide have been saved.")

        def load_slide(event):
            selected_indices = self.ppt_slide_listbox.curselection()
            if not selected_indices:
                return
            
            index = selected_indices[0]
            slide_data = self.ppt_slides[index]
            
            self.ppt_slide_title.delete(0, END)
            self.ppt_slide_title.insert(0, slide_data.get("title", ""))
            
            self.ppt_slide_content.delete("1.0", END)
            self.ppt_slide_content.insert("1.0", slide_data.get("content", ""))

        self.ppt_slide_listbox.bind('<<ListboxSelect>>', load_slide)

        def save_ppt():
            if not self.ppt_slides:
                messagebox.showwarning("No Slides", "Add at least one slide first")
                return
            
            filename = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint","*.pptx")])
            if filename:
                prs = Presentation()
                for s in self.ppt_slides:
                    layout = prs.slide_layouts[1] # Title and Content layout
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = s.get("title","")
                    slide.placeholders[1].text = s.get("content","")
                prs.save(filename)
                messagebox.showinfo("Saved", f"PowerPoint saved: {filename}")

        def upload_ppt():
            filename = filedialog.askopenfilename(filetypes=[("PowerPoint","*.pptx")])
            if filename:
                prs = Presentation(filename)
                self.ppt_slides.clear()
                self.ppt_slide_listbox.delete(0, END)
                
                for i, slide in enumerate(prs.slides):
                    title = slide.shapes.title.text if slide.shapes.title else ""
                    content = ""
                    if slide.placeholders and len(slide.placecolors) > 1:
                        content = slide.placeholders[1].text
                    else: # Fallback for non-standard layouts
                        for shp in slide.shapes:
                            if shp.has_text_frame and shp != slide.shapes.title:
                                content += shp.text + "\n"
                                
                    self.ppt_slides.append({"title": title, "content": content})
                    self.ppt_slide_listbox.insert(END, f"{i + 1}: {title}")
                
                clear_ppt_fields()
                messagebox.showinfo("Uploaded", f"{len(self.ppt_slides)} slides loaded from file.")

        # --- Button Frame ---
        ppt_btn_frame = Frame(win)
        ppt_btn_frame.pack(pady=5)
        
        Button(ppt_btn_frame, text="Add as New Slide", command=add_slide).pack(side=LEFT, padx=5)
        Button(ppt_btn_frame, text="Save Edit to Selected", command=save_slide_edit).pack(side=LEFT, padx=5)
        Button(ppt_btn_frame, text="Clear Fields", command=clear_ppt_fields).pack(side=LEFT, padx=5)
        
        Button(ppt_btn_frame, text="Upload PowerPoint", command=upload_ppt).pack(side=LEFT, padx=5)
        Button(ppt_btn_frame, text="Save PowerPoint", command=save_ppt).pack(side=LEFT, padx=5)

    # === PDF Merge Editor ===
    def open_pdf_window(self):
        win = Toplevel(self.root)
        win.title("PDF Merge")
        win.geometry("500x350") # Increased height for new button

        Label(win, text="Select PDFs to merge", font=("Arial",14)).pack(pady=10)
        listbox = Listbox(win, selectmode=MULTIPLE)
        listbox.pack(fill=BOTH, expand=True, padx=10, pady=5)

        def add_files():
            files = filedialog.askopenfilenames(filetypes=[("PDF Files","*.pdf")])
            for f in files:
                listbox.insert(END, f)

        def remove_selected():
            selected_indices = listbox.curselection()
            # Iterate backwards to avoid index shifting
            for i in reversed(selected_indices):
                listbox.delete(i)

        def merge_pdfs():
            files = listbox.get(0, END)
            if not files:
                messagebox.showwarning("No Files", "No PDF files selected")
                return
            
            output_file = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF","*.pdf")])
            if not output_file:
                return
            
            pdf_writer = PdfWriter()
            try:
                for f in files:
                    reader = PdfReader(f)
                    for page in reader.pages:
                        pdf_writer.add_page(page)
                
                with open(output_file, "wb") as f_out:
                    pdf_writer.write(f_out)
                
                messagebox.showinfo("Merged", f"PDF merged: {output_file} ({len(files)} files)")
                listbox.delete(0, END) # Clear list on success
            except Exception as e:
                messagebox.showerror("Error", f"An error occurred during merging: {e}")

        btn_frame_pdf = Frame(win)
        btn_frame_pdf.pack(pady=5)

        Button(btn_frame_pdf, text="Add PDF Files", command=add_files).pack(side=LEFT, padx=5)
        Button(btn_frame_pdf, text="Remove Selected", command=remove_selected).pack(side=LEFT, padx=5)
        Button(btn_frame_pdf, text="Merge PDFs", command=merge_pdfs).pack(side=LEFT, padx=5)


# === Run App ===
if __name__ == "__main__":
    root = Tk()
    app = OfficeSuiteApp(root)
    root.mainloop()